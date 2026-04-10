"""
Generate a professional course PPTX for LEAD GEN x AI Powered Tool.
Dark theme, green accent - matches the product branding.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Branding colors ──
BG       = RGBColor(15, 15, 15)
GREEN    = RGBColor(34, 197, 94)
WHITE    = RGBColor(255, 255, 255)
GRAY     = RGBColor(160, 160, 160)
DARK_CARD = RGBColor(28, 28, 28)
YELLOW   = RGBColor(250, 204, 21)
RED      = RGBColor(239, 68, 68)
LIGHT    = RGBColor(220, 220, 220)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def tf_setup(shape, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet_text(tf, text, font_size=14, color=LIGHT, bold=False, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_before = Pt(4)
    return p


def green_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()


def footer_text(slide, text="LEAD GEN x AI Powered Tool Course"):
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(5), Inches(0.4),
                 text, font_size=9, color=GRAY)


def section_number(slide, num):
    shape = add_shape(slide, Inches(0.5), Inches(0.3), Inches(0.6), Inches(0.45), GREEN)
    tf_setup(shape, str(num), font_size=18, color=BG, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
#  SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
green_bar(slide)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "LEAD GEN x AI Powered Tool", font_size=44, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
             "Complete Course: Scrape Leads, Build Websites, Pitch Clients", 
             font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.8),
             "Turn Google Maps data into a profitable web design business - using 100% free tools",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Badges
for i, txt in enumerate(["Free Tools Only", "No Coding Experience Needed", "Start Earning Today"]):
    shape = add_shape(slide, Inches(3 + i * 2.6), Inches(5.6), Inches(2.4), Inches(0.5), DARK_CARD)
    tf_setup(shape, txt, font_size=11, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

footer_text(slide)


# ═══════════════════════════════════════════════════════
#  SLIDE 2: COURSE AGENDA
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Course Agenda", font_size=32, color=GREEN, bold=True)

modules = [
    ("Module 1", "Setup & Installation", "VS Code, Git, Node.js, GitHub, Vercel, Chrome"),
    ("Module 2", "LEAD GEN x AI Powered Tool Walkthrough", "Install extension, run first scrape, export CSV"),
    ("Module 3", "Lead Qualification", "Identify HOT leads, audit websites, prioritize targets"),
    ("Module 4", "Building Websites with AI", "Claude prompts, HTML/CSS generation, local preview"),
    ("Module 5", "Free Deployment to Vercel", "Git push, auto-deploy, live demo URL"),
    ("Module 6", "Pitching Clients", "Cold email templates, follow-up sequences, closing deals"),
    ("Module 7", "Pricing & Packaging", "Service tiers, retainers, upsells"),
    ("Module 8", "Scaling Your Business", "Batch workflows, portfolio, repeat clients"),
]

for i, (mod, title, desc) in enumerate(modules):
    row = i // 2
    col = i % 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.5 + row * 1.4)
    
    card = add_shape(slide, x, y, Inches(5.8), Inches(1.2), DARK_CARD)
    # Module number
    num_shape = add_shape(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.9), Inches(0.35), GREEN)
    tf_setup(num_shape, mod, font_size=10, color=BG, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x + Inches(1.2), y + Inches(0.1), Inches(4.3), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True)
    # Description
    add_text_box(slide, x + Inches(1.2), y + Inches(0.55), Inches(4.3), Inches(0.5),
                 desc, font_size=11, color=GRAY)

footer_text(slide)


# ═══════════════════════════════════════════════════════
#  SLIDE 3: WHAT YOU'LL BUILD
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "What You'll Build By The End", font_size=32, color=GREEN, bold=True)

outcomes = [
    ("Scrape 100+ Leads", "Extract business names, emails, phones,\nratings, websites from Google Maps in minutes"),
    ("Build Free Websites", "Use Claude AI to generate complete,\nprofessional websites - zero coding needed"),
    ("Deploy for Free", "Push to Vercel and get a live URL\nin under 60 seconds"),
    ("Pitch & Close Clients", "Send personalized cold emails with\nlive demo links - get paying clients"),
]

for i, (title, desc) in enumerate(outcomes):
    x = Inches(0.5 + i * 3.1)
    card = add_shape(slide, x, Inches(1.8), Inches(2.9), Inches(3.5), DARK_CARD)
    
    # Number circle
    num = add_shape(slide, x + Inches(1.05), Inches(2.1), Inches(0.6), Inches(0.6), GREEN)
    tf_setup(num, str(i + 1), font_size=22, color=BG, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.5), Inches(2.5), Inches(1.2),
                 desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(0.5),
             "Total investment: $0  |  All tools are free  |  Start earning from Day 1",
             font_size=14, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 1: SETUP & INSTALLATION
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 1)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "Module 1: Setup & Installation", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Install all free tools needed for the entire course", font_size=14, color=GRAY)

tools = [
    ("VS Code", "Free code editor by Microsoft", "code.visualstudio.com", "Write & edit website code"),
    ("Git", "Version control system", "git-scm.com", "Track changes, push to GitHub"),
    ("Node.js", "JavaScript runtime (LTS)", "nodejs.org", "Required by some tools & Vercel"),
    ("GitHub", "Free code hosting", "github.com", "Store your code online"),
    ("Vercel", "Free web hosting", "vercel.com", "Deploy websites instantly"),
    ("Chrome", "Browser + extensions", "google.com/chrome", "Run LEAD GEN x AI Powered Tool"),
]

for i, (name, desc, url, purpose) in enumerate(tools):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.7 + row * 2.6)
    
    card = add_shape(slide, x, y, Inches(3.9), Inches(2.3), DARK_CARD)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4),
                 name, font_size=18, color=GREEN, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.5), Inches(0.3),
                 desc, font_size=12, color=WHITE)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.5), Inches(0.3),
                 url, font_size=11, color=YELLOW)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.4), Inches(3.5), Inches(0.6),
                 purpose, font_size=11, color=GRAY)

footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 1: INSTALLATION STEPS
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 1)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "Installation Steps (Demo)", font_size=30, color=WHITE, bold=True)

steps = [
    "1. Download & install VS Code - check 'Add to PATH'",
    "2. Install extensions: Live Server, Prettier, HTML CSS Support",
    "3. Download & install Git - select VS Code as default editor",
    "4. Download & install Node.js LTS - check 'Add to PATH'",
    "5. Create free GitHub account at github.com/signup",
    "6. Create free Vercel account - sign up with GitHub",
    "7. Install LEAD GEN x AI Powered Tool extension in Chrome",
    "",
    "Verify in VS Code terminal (Ctrl + `):",
    "   code --version",
    "   git --version", 
    "   node --version",
    "   npm --version",
]

tf = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5.5),
                  steps[0], font_size=15, color=LIGHT)
for s in steps[1:]:
    if s == "":
        add_bullet_text(tf, "", font_size=8, color=BG)
    elif s.startswith("   "):
        add_bullet_text(tf, s, font_size=14, color=GREEN, bold=True)
    elif s.startswith("Verify"):
        add_bullet_text(tf, s, font_size=15, color=YELLOW, bold=True)
    else:
        add_bullet_text(tf, s, font_size=15, color=LIGHT)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(10), Inches(0.4),
             "LIVE DEMO: Walk through each installation step together",
             font_size=13, color=GREEN, bold=True)
footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 2: LEAD GEN x AI POWERED TOOL
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 2)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "Module 2: LEAD GEN x AI Powered Tool Walkthrough", font_size=30, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
             "How It Works", font_size=20, color=GREEN, bold=True)

how_steps = [
    "1. Install the Chrome extension (free credits included)",
    "2. Click the extension icon in Chrome toolbar",
    "3. Type any search: 'dentists in Miami'",
    "4. Set how many results you want (e.g., 50)",
    "5. Click 'Start Scraping' - sit back and watch",
    "6. Extension auto-extracts: name, phone, email,\n    address, rating, reviews, website, hours",
    "7. Results auto-export to CSV + open in viewer",
]

tf = add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                  how_steps[0], font_size=15, color=LIGHT)
for s in how_steps[1:]:
    add_bullet_text(tf, s, font_size=15, color=LIGHT)

# Right side - what you get
card = add_shape(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(5), DARK_CARD)
add_text_box(slide, Inches(7.3), Inches(1.4), Inches(5), Inches(0.4),
             "Data You Get Per Lead:", font_size=16, color=YELLOW, bold=True)

fields = ["Business Name", "Phone Number", "Email Address", "Full Address",
          "Google Rating", "Review Count", "Website URL", "Business Hours",
          "Social Media Links"]
tf2 = add_text_box(slide, Inches(7.3), Inches(2.0), Inches(5), Inches(3.5),
                   fields[0], font_size=14, color=GREEN)
for f in fields[1:]:
    add_bullet_text(tf2, f, font_size=14, color=GREEN)

add_text_box(slide, Inches(7.3), Inches(5.5), Inches(5), Inches(0.4),
             "All in one click. No coding required.", font_size=12, color=GRAY)
footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 2: LIVE DEMO
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 2)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "Live Demo: First Scrape", font_size=30, color=WHITE, bold=True)

card = add_shape(slide, Inches(1), Inches(1.5), Inches(11), Inches(4.5), DARK_CARD)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(0.5),
             "DEMO: Scraping 'plumbers in Chicago'", font_size=20, color=GREEN, bold=True)

demo_steps = [
    "1. Open Chrome > Click LEAD GEN x AI Powered Tool extension icon",
    "2. Type: plumbers in Chicago",
    "3. Set max results: 20",
    "4. Click 'Start Scraping'",
    "5. Watch it auto-navigate Google Maps and extract each listing",
    "6. Results appear in real-time in the popup",
    "7. CSV auto-downloads when done",
    "8. Open the built-in viewer to browse results",
]

tf = add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(3),
                  demo_steps[0], font_size=15, color=LIGHT)
for s in demo_steps[1:]:
    add_bullet_text(tf, s, font_size=15, color=LIGHT)

add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
             "[ SCREEN RECORDING: Show the full scraping process live ]",
             font_size=14, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 3: LEAD QUALIFICATION
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 3)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "Module 3: Lead Qualification", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Not all leads are equal. Here's how to find the money leads.", font_size=14, color=GRAY)

# HOT leads
hot = add_shape(slide, Inches(0.5), Inches(1.7), Inches(3.8), Inches(4.5), DARK_CARD)
hot_badge = add_shape(slide, Inches(0.7), Inches(1.9), Inches(1), Inches(0.4), RGBColor(239, 68, 68))
tf_setup(hot_badge, "HOT", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.85), Inches(1.85), Inches(2), Inches(0.4),
             "Highest Priority", font_size=14, color=RED, bold=True)
hot_items = [
    "No website at all",
    "4+ star rating, 10+ reviews",
    "High-value niche (dental, legal,\nplumbing, real estate, clinics)",
    "Active business (recent reviews)",
    "Has email or phone to contact",
]
tf = add_text_box(slide, Inches(0.7), Inches(2.5), Inches(3.3), Inches(3.5),
                  hot_items[0], font_size=13, color=LIGHT)
for item in hot_items[1:]:
    add_bullet_text(tf, item, font_size=13, color=LIGHT)

# WARM leads
warm = add_shape(slide, Inches(4.7), Inches(1.7), Inches(3.8), Inches(4.5), DARK_CARD)
warm_badge = add_shape(slide, Inches(4.9), Inches(1.9), Inches(1.2), Inches(0.4), YELLOW)
tf_setup(warm_badge, "WARM", font_size=14, color=BG, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.25), Inches(1.85), Inches(2), Inches(0.4),
             "Follow Up", font_size=14, color=YELLOW, bold=True)
warm_items = [
    "Has outdated/ugly website",
    "Not mobile-friendly",
    "No SSL (http instead of https)",
    "Missing contact form or CTA",
    "Good reviews but poor web presence",
]
tf = add_text_box(slide, Inches(4.9), Inches(2.5), Inches(3.3), Inches(3.5),
                  warm_items[0], font_size=13, color=LIGHT)
for item in warm_items[1:]:
    add_bullet_text(tf, item, font_size=13, color=LIGHT)

# COLD leads
cold = add_shape(slide, Inches(8.9), Inches(1.7), Inches(3.8), Inches(4.5), DARK_CARD)
cold_badge = add_shape(slide, Inches(9.1), Inches(1.9), Inches(1.2), Inches(0.4), GRAY)
tf_setup(cold_badge, "COLD", font_size=14, color=BG, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.45), Inches(1.85), Inches(2), Inches(0.4),
             "Skip These", font_size=14, color=GRAY, bold=True)
cold_items = [
    "Already has a good website",
    "Low ratings (below 3 stars)",
    "Very few reviews (< 5)",
    "Chain/franchise (corporate sites)",
    "No contact info available",
]
tf = add_text_box(slide, Inches(9.1), Inches(2.5), Inches(3.3), Inches(3.5),
                  cold_items[0], font_size=13, color=LIGHT)
for item in cold_items[1:]:
    add_bullet_text(tf, item, font_size=13, color=LIGHT)

add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
             "Pro Tip: Use the AI Bulk Qualification Prompt to score 50+ leads at once",
             font_size=13, color=GREEN, bold=True)
footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 4: BUILDING WEBSITES WITH AI
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 4)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "Module 4: Building Websites with AI", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Use Claude (free) to generate complete websites from your scraped data", font_size=14, color=GRAY)

# Workflow
workflow_steps = [
    ("Pick a HOT Lead", "Choose a business with no website\nor a terrible one from your CSV"),
    ("Fill the Prompt", "Copy the Website Build prompt,\nfill in placeholders from CSV data"),
    ("Paste into Claude", "Go to claude.ai (free),\npaste the prompt, get full code"),
    ("Preview Locally", "Save files in VS Code,\nright-click > Open with Live Server"),
]

for i, (title, desc) in enumerate(workflow_steps):
    x = Inches(0.5 + i * 3.15)
    card = add_shape(slide, x, Inches(1.7), Inches(2.9), Inches(2.5), DARK_CARD)
    
    num = add_shape(slide, x + Inches(0.15), Inches(1.85), Inches(0.45), Inches(0.45), GREEN)
    tf_setup(num, str(i + 1), font_size=18, color=BG, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.7), Inches(1.85), Inches(2), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.5), Inches(1.2),
                 desc, font_size=12, color=GRAY)

# What Claude generates
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(0.4),
             "What Claude Generates For You:", font_size=18, color=GREEN, bold=True)

gen_items = [
    ("Hero Section", "With business name, tagline, and Call-to-Action button"),
    ("Services Section", "Tailored to the business niche"),
    ("Testimonials", "Using their actual Google reviews"),
    ("Contact Form", "Working form via Formspree (free)"),
    ("Google Map Embed", "Showing their exact location"),
    ("Full SEO", "Meta tags, schema markup, Open Graph"),
]

for i, (title, desc) in enumerate(gen_items):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(5.0 + row * 0.7)
    add_text_box(slide, x, y, Inches(4), Inches(0.3),
                 f"{title}: {desc}", font_size=12, color=LIGHT)

footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 4: AI PROMPTS (reference)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 4)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "AI Prompt Library (Included)", font_size=30, color=WHITE, bold=True)

prompts = [
    ("1. Full Website Build", "Generates complete HTML/CSS website\nfrom scraped lead data", GREEN),
    ("2. Website Audit + Emails", "Audits their site + creates 3 cold\nemail templates + follow-up sequence", GREEN),
    ("3. Landing Page Generator", "High-converting single page with\nCTA, reviews, FAQ sections", GREEN),
    ("4. Bulk Lead Qualifier", "Scores 50+ leads as HOT/WARM/COLD\nfrom your CSV export", GREEN),
    ("5. Pitch Deck Generator", "Full pitch document with pricing,\nROI calculation, comparisons", GREEN),
    ("6. System Prompt", "Turn Claude into your full-stack\nassistant for the entire workflow", YELLOW),
]

for i, (title, desc, color) in enumerate(prompts):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 1.8)
    
    card = add_shape(slide, x, y, Inches(5.9), Inches(1.6), DARK_CARD)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.3), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), Inches(5.3), Inches(0.8),
                 desc, font_size=12, color=GRAY)

add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
             "All prompts included in the AI_Prompts_and_Pitch_Templates.pdf",
             font_size=13, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 5: DEPLOYMENT
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 5)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "Module 5: Free Deployment to Vercel", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Get a live URL in under 60 seconds - completely free", font_size=14, color=GRAY)

deploy_steps = [
    ("Create GitHub Repo", "Go to github.com/new\nName it: client-business-name\nSelect Public, add README"),
    ("Clone to Your Computer", "In VS Code terminal:\ngit clone https://github.com/\n  YOU/client-name.git\ncd client-name"),
    ("Add Website Files", "Paste the AI-generated code\ninto index.html and styles.css\nSave all files"),
    ("Push to GitHub", "git add .\ngit commit -m 'Initial website'\ngit push"),
    ("Import in Vercel", "Go to vercel.com/new\nImport the GitHub repo\nClick Deploy"),
    ("Live in 30 Seconds!", "Your site is live at:\nclient-name.vercel.app\nShare this link in pitch email"),
]

for i, (title, desc) in enumerate(deploy_steps):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.7 + row * 2.7)
    
    card = add_shape(slide, x, y, Inches(3.9), Inches(2.4), DARK_CARD)
    num = add_shape(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.45), Inches(0.45), GREEN)
    tf_setup(num, str(i + 1), font_size=18, color=BG, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.7), y + Inches(0.15), Inches(3), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(1.5),
                 desc, font_size=12, color=GRAY)

footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 6: PITCHING CLIENTS
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 6)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "Module 6: Pitching Clients", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "3 proven email templates + follow-up sequence", font_size=14, color=GRAY)

emails = [
    ("Email A", '"I Built You a Website"', "HIGHEST CONVERSION",
     "Build the site FIRST, send the live link.\nThey see immediate value before replying.\nInclude their business name + specific details.",
     RGBColor(239, 68, 68)),
    ("Email B", '"Quick Website Audit"', "CONSULTATIVE",
     "Point out 2-3 specific issues with their\ncurrent site (slow, not mobile-friendly).\nOffer a free redesign preview.",
     YELLOW),
    ("Email C", '"Fellow Local Business"', "RELATIONSHIP",
     "Position as a local professional.\nMention their strong Google reviews.\nSoft CTA: quick call or reply.",
     GREEN),
]

for i, (label, subject, tag, desc, color) in enumerate(emails):
    x = Inches(0.5 + i * 4.15)
    card = add_shape(slide, x, Inches(1.7), Inches(3.9), Inches(3.8), DARK_CARD)
    
    badge = add_shape(slide, x + Inches(0.15), Inches(1.85), Inches(1.4), Inches(0.4), color)
    tf_setup(badge, f"  {tag}", font_size=10, color=BG, bold=True)
    
    add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(3.5), Inches(0.4),
                 f"{label}: {subject}", font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(3.5), Inches(2),
                 desc, font_size=12, color=GRAY)

# Follow-up
add_text_box(slide, Inches(0.8), Inches(5.8), Inches(10), Inches(0.4),
             "Follow-Up Sequence:", font_size=16, color=GREEN, bold=True)
fu = add_text_box(slide, Inches(0.8), Inches(6.3), Inches(10), Inches(0.8),
                  "Day 3: Quick check-in (did you see the demo?)", font_size=13, color=LIGHT)
add_bullet_text(fu, "Day 7: Add social proof (other local businesses you've helped)", font_size=13, color=LIGHT)
add_bullet_text(fu, "Day 14: Final follow-up with limited-time offer", font_size=13, color=LIGHT)
footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 7: PRICING
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 7)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "Module 7: Pricing & Packaging", font_size=30, color=WHITE, bold=True)

tiers = [
    ("STARTER", "$200 - $500", [
        "Single landing page",
        "Mobile responsive",
        "Contact form",
        "Basic SEO",
        "Deploy to Vercel",
    ], GRAY),
    ("STANDARD", "$500 - $1,500", [
        "Multi-page website (3-5 pages)",
        "Contact form + Google Maps",
        "Full SEO + schema markup",
        "Google Business optimization",
        "Social media links",
    ], GREEN),
    ("PREMIUM", "$1,500 - $3,000+", [
        "Everything in Standard",
        "Custom design / branding",
        "Blog or portfolio section",
        "Analytics setup",
        "3 months maintenance included",
    ], YELLOW),
]

for i, (name, price, features, color) in enumerate(tiers):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(1.5), Inches(3.7), Inches(4.5), DARK_CARD)
    
    if i == 1:  # highlight standard
        border_shape = add_shape(slide, x - Inches(0.03), Inches(1.47), Inches(3.76), Inches(4.56), GREEN)
        border_shape.fill.background()
        border_shape.line.color.rgb = GREEN
        border_shape.line.width = Pt(2)
    
    badge = add_shape(slide, x + Inches(0.15), Inches(1.7), Inches(1.5), Inches(0.4), color)
    tf_setup(badge, f"  {name}", font_size=12, color=BG, bold=True)
    
    add_text_box(slide, x + Inches(0.2), Inches(2.3), Inches(3.3), Inches(0.5),
                 price, font_size=24, color=WHITE, bold=True)
    
    tf = add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(3.3), Inches(2.5),
                      features[0], font_size=12, color=LIGHT)
    for f in features[1:]:
        add_bullet_text(tf, f, font_size=12, color=LIGHT)

# Retainer
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
             "Recurring Revenue: Offer monthly maintenance at $50-$150/month (hosting, updates, backups)",
             font_size=14, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
footer_text(slide)


# ═══════════════════════════════════════════════════════
#  MODULE 8: SCALING
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)
section_number(slide, 8)

add_text_box(slide, Inches(1.3), Inches(0.25), Inches(8), Inches(0.7),
             "Module 8: Scaling Your Business", font_size=30, color=WHITE, bold=True)

scale_items = [
    ("Batch Processing", "Scrape 50-100 leads per niche per city.\nQualify in bulk with AI. Build 3-5 demo\nsites per week."),
    ("Build a Portfolio", "Every client site becomes portfolio proof.\nScreenshot before/after. Add testimonials\nfrom happy clients."),
    ("Niche Down", "Focus on 2-3 high-value niches.\nBecome the 'dentist website guy' or\n'restaurant web expert' in your area."),
    ("Repeat Clients", "Offer monthly maintenance retainers.\nUpsell Google Business optimization.\nReferral discounts for introductions."),
]

for i, (title, desc) in enumerate(scale_items):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.5 + row * 2.5)
    
    card = add_shape(slide, x, y, Inches(5.9), Inches(2.2), DARK_CARD)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(5), Inches(0.4),
                 title, font_size=16, color=GREEN, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), Inches(5.3), Inches(1.3),
                 desc, font_size=13, color=LIGHT)

add_text_box(slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             "Goal: 5 clients x $500 avg = $2,500/month  |  10 clients x $100/mo retainer = $1,000/mo recurring",
             font_size=14, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
footer_text(slide)


# ═══════════════════════════════════════════════════════
#  FINAL SLIDE: COMPLETE WORKFLOW
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Your Complete Workflow", font_size=32, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

flow = [
    ("SCRAPE", "Extract leads\nfrom Google Maps"),
    ("QUALIFY", "Score leads as\nHOT / WARM / COLD"),
    ("BUILD", "Generate website\nwith Claude AI"),
    ("DEPLOY", "Push to Vercel\nget live URL"),
    ("PITCH", "Send cold email\nwith demo link"),
    ("CLOSE", "Deliver site\ncollect payment"),
]

for i, (title, desc) in enumerate(flow):
    x = Inches(0.3 + i * 2.15)
    card = add_shape(slide, x, Inches(1.8), Inches(1.9), Inches(2.5), DARK_CARD)
    
    num = add_shape(slide, x + Inches(0.65), Inches(2.0), Inches(0.5), Inches(0.5), GREEN)
    tf_setup(num, str(i + 1), font_size=18, color=BG, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.1), Inches(2.7), Inches(1.7), Inches(0.3),
                 title, font_size=13, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(3.1), Inches(1.7), Inches(0.8),
                 desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# Resources
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(10), Inches(0.4),
             "Resources Included:", font_size=18, color=WHITE, bold=True)

resources = [
    "AI_Prompts_and_Pitch_Templates.pdf - All 6 prompts ready to use",
    "Beginner_Setup_Tutorial.pdf - Step-by-step installation guide",
    "LEAD GEN x AI Powered Tool Extension - Chrome extension with free credits",
]
tf = add_text_box(slide, Inches(0.8), Inches(5.4), Inches(10), Inches(1.5),
                  resources[0], font_size=14, color=GREEN)
for r in resources[1:]:
    add_bullet_text(tf, r, font_size=14, color=GREEN)

add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
             "Total cost to start: $0  |  Potential earnings: $2,500+/month",
             font_size=16, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
footer_text(slide)


# ═══════════════════════════════════════════════════════
#  THANK YOU SLIDE
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
green_bar(slide)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
             "Now Go Get Your First Client!", font_size=40, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             "Scrape  >  Build  >  Deploy  >  Pitch  >  Get Paid", 
             font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "map-scrapper-five.vercel.app", 
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
             "Questions? Drop them in the comments!", 
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'LEAD_GEN_x_AI_Powered_Tool_Course.pptx')
prs.save(output_path)
print(f'[OK] {output_path}')
